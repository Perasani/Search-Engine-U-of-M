import os, errno #import library
import math#import library
import time#import library#import library
import operator #import library
import collections #import library
from collections import OrderedDict #import library
from collections import deque #import library
import queue #import library
import shutil #import library 
import pickle #import library
import re #import library
import urllib #import library
from urllib.parse import urlsplit #import library
from urllib.parse import urljoin #import library
import requests #import library
from bs4 import BeautifulSoup, Tag #import library
from nltk.stem.porter import PorterStemmer #import library
import datetime #import library
import numpy as np #import library
import docx2txt #import library
from pptx import Presentation #import library
import copy #import library

stopwords = {} #declare global variables
term_doc_freq_vector = {} #declare global variables
doc_term_freq_vector = {} #declare global variables
doc_term_freq_vector_norm = {} #declare global variables
page_doc_map = {} #declare global variables
doc_page_map = {} #declare global variables
page_ref_count = {} #declare global variables
doc_count = 0 #declare global variables
link_queue = queue.Queue() #declare global variables
last_doc_index = -1 #declare global variables
page_queued_map = {} #declare global variables
start_time = time.time() #declare global variables
total_number_docs = 0 #declare global variables
crawled_web_dir = "web_text_crawled" #declare global variables
crawled_web_dir_conv_need = "web_docs_crawled" #declare global variables
crawled_web_dir_preprocessed = "web_text_preprocessed" #declare global variables
output_web_dir = "output" #declare global variables
stopword_path = "english.stopwords.txt" #declare global variables
list_dir = [crawled_web_dir, crawled_web_dir_conv_need, crawled_web_dir_preprocessed] #declare global variables
url = "http://www.cs.memphis.edu/~vrus/teaching/ir-websearch/" #declare url
domain = "memphis.edu" #declare domain
total_number_docs = 10000 #declare total number of webpages
def create_directory(directory): #create directory 
    try:
        os.makedirs(directory)  #create directory
    except OSError as e:  #create directory
        if(e.errno != errno.EEXIST):  #create directory
            raise
    pass 
def create_directories(list_dir):
    for dir_i in list_dir:
        print(dir_i)  #create directory
        create_directory(dir_i) 
def delete_directory(dir_name):  #create directory
    if(os.path.isdir(dir_name)):
        try:  #create directory
            shutil.rmtree(dir_name)
        except OSError as e:
            if(e.errno != errno.EEXIST):
                raise
        pass   #create directory
def delete_directories(list_dir):
    for dir_i in list_dir:
        print(dir_i)
        delete_directory(dir_i)  #create directory
def delete_file(path):
    try:
        os.remove(path)
    except WindowsError:  #create directory
        print("failed deleting: " + path)
        pass  #create directory
def save_text(text, dir_path, file_name):
    text_file = open(dir_path+"\\"+file_name, "w") #create directory
    text_file.write(text)
        pass  #create directory
    text_file.close() 
def load_stopwords(filepath): #create directory
        pass  #create directory
    with open(filepath, 'r') as content_file:  
        for line in content_file:
            line = line.strip() #create directory
            stopwords[line] = 1 
def save_obj(obj, name, key_or_val, order):
    filename = name + ".p"
    if(key_or_val == "key" and order == "auto"): #create directory
        sorted_x = sorted(obj.items(), key=operator.itemgetter(0)) #create directory
    elif(key_or_val == "key" and order == "reverse"): #create directory
        sorted_x = sorted(obj.items(), key=operator.itemgetter(0), reverse=True) #create directory
    elif(key_or_val == "value" and order == "auto"):
        sorted_x = sorted(obj.items(), key=operator.itemgetter(1)) #create directory
    elif(key_or_val == "value" and order == "reverse"):
        sorted_x = sorted(obj.items(), key=operator.itemgetter(1), reverse=True)
    if (os.path.isfile(filename)): #create directory
        os.remove(filename)
    pickle.dump( obj, open( filename, "wb" ) ) #create directory
def save_obj_without_sort(obj, name):
    pickle.dump( obj, open( name + ".p", "wb" ) )
def save_obj_no_sort(obj, name): #declare global variables
    filename = name + ".p"
    
    if (os.path.isfile(filename)):
        os.remove(filename) #create directory
        
    pickle.dump( obj, open( filename, "wb" ) )
def save_obj_no_sort_w(queue1, filename):
    link_list = [] #create directory
    
    new_queue = queue.Queue() #create directory
    new_queue.queue = copy.deepcopy(queue1.queue)
    if (os.path.isfile(filename)):
        os.remove(filename)
         #create directory
    while(new_queue.empty() == False): #create directory
        link_list.append(new_queue.get())
        
    filename = filename + ".p"
     #create directory
    pickle.dump( link_list, open( filename, "wb" ) ) #create directory
def load_obj(name):
    file = open(name,'rb') #declare global variables
    object_file = pickle.load(file)
    file.close() #create directory
    return object_file
def load_obj_no_sort(name): #create directory
    file = open(name,'rb')
    object_file = pickle.load(file)
    file.close()
    
    return object_file #create directory
def load_obj_no_sort_w(filename):
    global link_queue
    link_list =[]
    #create directory
    if (os.path.isfile(filename)):
        file = open(filename,'rb')
        link_list = pickle.load(file)
        file.close()
        
        for link in link_list:
            link_queue.put(link)
        return link_queue
    
    else:
        print("no file found")
        return
def reset_global_variables():
    global stopwords #declare global variables
    global term_doc_freq_vector #declare global variables
    global doc_term_freq_vector  #declare global variables
    global page_doc_map  #declare global variables
    global doc_page_map #declare global variables
    global page_ref_count #declare global variables
    global doc_count #declare global variables
    global link_queue #declare global variables
    global last_doc_index #declare global variables
    global page_queued_map #declare global variables
    global start_time #declare global variables

    stopwords = {}
    term_doc_freq_vector = {}
    doc_term_freq_vector = {}
    page_doc_map = {}
    doc_page_map = {}
    page_ref_count = {}
    doc_count = 0
    link_queue = queue.Queue()
    last_doc_index = -1
    page_queued_map = {}
    start_time = time.time()
def delete_all_files():
    try:
        delete_directories(list_dir)
    except:
        pass
    try: 
        delete_file("doc_count.p")#declare global variables
    except:
        pass
    try: #declare global variables
        delete_file("doc_term_freq_vector.p")
    except: #declare global variables
        pass
    try: #declare global variables
        delete_file("doc_term_freq_vector_norm.p")
    except:
        pass
    try:
        delete_file("doc_url_map.p")
    except: #declare global variables
        pass
    try:
        delete_file("link_queue.p")
    except: #declare global variables
        pass
    try:
        delete_file("page_ref_count.p")
    except: #declare global variables
        pass
    try: #declare global variables
        delete_file("term_doc_freq_vector.p")
    except:
        pass
    try:
        delete_file("url_doc_map.p") #declare global variables
    except:
        pass
def save_all_obj():
    global page_doc_map #declare global variables
    global doc_page_map #declare global variables
    global page_ref_count #declare global variables
    global doc_count #declare global variables
    global link_queue #declare global variables
    
    save_obj(page_doc_map, "url_doc_map", "value", "auto") #declare global variables
    save_obj(doc_page_map, "doc_url_map", "key", "auto") #declare global variables
    save_obj_no_sort(doc_count, "doc_count") #declare global variables
    save_obj_no_sort_w(link_queue, "link_queue") #declare global variables
    save_obj(page_ref_count, "page_ref_count", "value", "reverse") #declare global variables
    print("saved objects")  
def save_all_obj_tfidf(doc_term_freq_vector_norm_new): #declare global variables
    global term_doc_freq_vector
    global doc_term_freq_vector #declare global variables
     #declare global variables
    save_obj_without_sort(term_doc_freq_vector, "term_doc_freq_vector")
    save_obj_without_sort(doc_term_freq_vector, "doc_term_freq_vector")
    save_obj_without_sort(doc_term_freq_vector_norm_new, "doc_term_freq_vector_norm") 
def load_all_obj():
    global page_doc_map
    global doc_page_map
    global page_ref_count
    global doc_count
    global link_queue
     #declare global variables
    page_doc_map = load_obj("url_doc_map.p")
    doc_page_map = load_obj("doc_url_map.p") #declare global variables
    doc_count = max(doc_page_map.keys()) #declare global variables
    link_queue = load_obj_no_sort_w("link_queue.p") #declare global variables
    page_ref_count = load_obj("page_ref_count.p")
    print("loaded objects")
def load_obj_search():
    global total_number_docs #declare global variables
    global doc_url_map
    global term_doc_freq_vector #declare global variables
    global doc_term_freq_vector
    global doc_term_freq_vector_norm #declare global variables
    
    stopword_path = "english.stopwords.txt"
    doc_url_map_file = "doc_url_map.p" #declare global variables
    term_doc_freq_file = "term_doc_freq_vector.p"
    doc_term_freq_file = "doc_term_freq_vector.p"
    doc_term_freq_file_norm = "doc_term_freq_vector_norm.p"
    doc_count = "doc_count.p" #declare global variables
    
    total_number_docs = load_obj(doc_count)
    doc_url_map = load_obj(doc_url_map_file)
    term_doc_freq_vector = load_obj(term_doc_freq_file) #declare global variables
    doc_term_freq_vector = load_obj(doc_term_freq_file)
    doc_term_freq_vector_norm = load_obj(doc_term_freq_file_norm)
    load_stopwords(stopword_path)
def format_time(start_time, end_time): #declare global variables
    elsapsed_time = end_time - start_time
    hr = int(elsapsed_time)//3600
    min_ = (int(elsapsed_time) - (hr * 3600))/60
    sec = int(elsapsed_time) - hr * 3600 - min_ * 60
    print("HH:Min:Sec > " + str(hr) +" hr " + str(min_) + " min "+ str(sec) + "sec")
def remove_url_frag_id(url):
    if ".php" in url:
        url = url.split('.php') #declare global variables
        if(len(url)>1):
            url =url[0] + ".php"
    elif ".aspx" in url: #declare global variables
        url = url.split('.aspx')
        if(len(url)>1):#declare global variables
            url =url[0] + ".aspx"   
    url = url.split('#')[0]
    
    return url
def remove_url_frag_simple(url):    #declare global variables
    url = url.split('#')[0]
    return url #declare global variables
def remove_slash_before_or_after(url, type_r):
    if(type_r == "before"): #declare global variables
        if url.startswith("/"):
            url = url[1:]
        return url #declare global variables
    
    elif(type_r == "after"):   
        if url[-1]=="/": #declare global variables
            url = url.rsplit('/', 1)[0]
        return url 
def strip_http_s(url):     #declare global variables
    url = url.replace("https://","")
    url = url.replace("http://","") #declare global variables
    url = url.rstrip('\/') 
    url = "http://"+ url
    
    return url#declare global variables
def check_if_in_domain(url, domain):
    if(domain in url):  #declare global variables
        return 1
    else:
        return 0#declare global variables
def is_excluded_type(extension):
    exclude_list = ["jpg", "jpeg", "png", "mp3", "mp4", "xlx"]
    if extension in exclude_list:#declare global variables
        return 1
    else:
        return 0
def check_valid_URL(url):
    url_reg = re.compile(
        r'^(?:http|ftp)s?://' # http:// or https://
        r'(?:(?:[A-Z0-9](?:[A-Z0-9-]{0,61}[A-Z0-9])?\.)+(?:[A-Z]{2,6}\.?|[A-Z0-9-]{2,}\.?)|'
        r'localhost|' #localhost...
        r'\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3})' # ...or ip#declare global variables
        r'(?::\d+)?' # optional port
        r'(?:/?|[/?]\S+)$', re.IGNORECASE)
    #declare global variables
    is_valid = url_reg.match(url)
    
    return is_valid
def get_page_extention(url):#declare global variables
    weblink_extention = url.rsplit('.', 1)[-1]#declare global variables
    return weblink_extention
def remove_hyper_link(text):
    URLless_string = re.sub(r'(?i)\b((?:https?://|www\d{0,3}[.]|[a-z0-9.\-]+[.][a-z]{2,4}     /)(?:[^\s()<>]+|\(([^\s()<>]+|(\([^\s()<>]+\)))*\))+(?:\(([^\s()<>]+|(\([^\s()<>]+\)))*\)|[^\s`!()\[\]{};:\'".,<>?«»“”‘’]))', '', text)
    return URLless_string
def remove_special_char(line):#declare global variables
    line = re.sub('[^a-zA-Z]+', ' ', line)
    return line     
def get_all_links(url, html):#declare global variables
    global domain#declare global variables
    global link_queue
    global page_queued_map
    #declare global variables#declare global variables
    soup = BeautifulSoup(html, "html.parser")#declare global variables
    links = soup.find_all('a')
    
    
    for tag in links:
        link = tag.get('href', None)
        #declare global variables
        if link is not None:
            try:
                link_extention = get_page_extention(link)#declare global variables
                
                if(link == "" or link == "#" or link_extention == "ppt" or is_excluded_type(link_extention) == 1):
                    a=1

                elif(link_extention in ["pdf", "docx", "pptx", "txt"]):
                    if(check_valid_URL(link)):#declare global variables
                        link_original = strip_http_s(link)
                        
                        if link_original not in page_queued_map:
                            page_queued_map[link_original] = 1#declare global variables
                            link_queue.put(link)

                    else:
                        modified_url = remove_url_frag_id(url)
                        modified_url = remove_slash_before_or_after(modified_url, "after")
                        modified_link = remove_slash_before_or_after(link, "before")
                        modified_link = modified_url + "/" + modified_link
                        
                        if(check_valid_URL(modified_link)):
                            link_original = strip_http_s(modified_link)
                            
                            if link_original not in page_queued_map:
                                page_queued_map[link_original] = 1
                                link_queue.put(modified_link)

                else:
                    is_valid = check_valid_URL(link)
                    if(is_valid):
                        modified_link = remove_url_frag_id(link)
                        modified_link = remove_slash_before_or_after(modified_link, "after")
                        link_original = strip_http_s(modified_link) #declare global variables
                        
                        if(check_if_in_domain(modified_link, domain) == 1):
                            
                            if link_original not in page_queued_map:
                                page_queued_map[link_original] = 1#declare global variables
                                link_queue.put(modified_link)

                    else: #declare global variables
                        modified_url= remove_url_frag_id(url)
                        modified_url = remove_slash_before_or_after(modified_url, "after")
                        modified_link = remove_slash_before_or_after(url, "before")
                        
                        if(modified_url!=modified_link):
                            
                            if  modified_url not in modified_link:
                                modified_link = modified_url + "/" + modified_link 
                                 #declare global variables
                                if(check_if_in_domain(modified_link, domain) == 1):
                                    if(check_valid_URL(modified_link)):
                                        link_original = strip_http_s(modified_link) #declare global variables
                                        
                                        if link_original not in page_queued_map:
                                            page_queued_map[link_original] = 1 #declare global variables
                                            link_queue.put(modified_link) 
                            else: 
                                if(check_if_in_domain(modified_link, domain) == 1):
                                    if(check_valid_URL(modified_link)):
                                        link_original = strip_http_s(modified_link)
                                        
                                        if link_original not in page_queued_map:
                                            page_queued_map[link_original] = 1
                                            link_queue.put(modified_link) #declare global variables
            except:
                continue
def pdf_to_text(input_pdf, file_name): #declare global variables
    global crawled_web_dir
    os.system(("pdftotext %s %s") %( input_pdf, crawled_web_dir+"//"+file_name)) #declare global variables
def pptx_to_text(book_path):
    prs = Presentation(book_path) #declare global variables
    text = "" #declare global variables
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = text_runs + run.text
    return text
def import_convert_preprocess(url, extension): #declare global variables
    global doc_count #declare global variables
    global crawled_web_dir_preprocessed
    global crawled_web_dir #declare global variables
    global crawled_web_dir_conv_need #declare global variables
    global page_doc_map  #declare global variables
    url_map_name = url #declare global variables
     #declare global variables
    if(url_map_name not in page_doc_map):
        page_doc_map[url_map_name] = -1
        page_ref_count[url_map_name] = 1
        
        try:
            doc_count_temp = doc_count + 1 #declare global variables
            book_name = "" #declare global variables
            if extension == "pdf":
                book_name = str(doc_count_temp) + ".pdf" #declare global variables
            elif extension == "docx":
                book_name = str(doc_count_temp) + ".docx"#declare global variables
            elif extension == "pptx":
                book_name = str(doc_count_temp) + ".pptx" #declare global variables

            book_path = crawled_web_dir_conv_need + "\\" + book_name
            
            a = requests.get(url, stream=True)
             #declare global variables
            with open(book_path, 'wb') as book:   
                for block in a.iter_content(512):
                    if not block:
                        break
                    book.write(block)
                    
            book.close()
#declare global variables
            file_name = str(doc_count_temp)+".txt"
            file_path = crawled_web_dir+ "\\" + file_name
            is_valid_for_indexing = 555
            if extension == "pdf": #declare global variables
                pdf_to_text(book_path, file_name)
                is_valid_for_indexing = preprocess_one_doc_from_pdf(crawled_web_dir, file_name, crawled_web_dir_preprocessed)
                
            elif extension == "docx": #declare global variables
                text = docx2txt.process(book_path)
                save_text(text, crawled_web_dir, file_name)
                is_valid_for_indexing = preprocess_one_doc(crawled_web_dir, file_name, crawled_web_dir_preprocessed)
                
            elif extension == "pptx": #declare global variables
                text = pptx_to_text(book_path)
                save_text(text, crawled_web_dir, file_name)
                is_valid_for_indexing = preprocess_one_doc(crawled_web_dir, file_name, crawled_web_dir_preprocessed)
        

            if(is_valid_for_indexing == 1) : #declare global variables
                doc_count = doc_count + 1
                page_doc_map[url_map_name] = doc_count
                doc_page_map[doc_count] = url_map_name
                page_ref_count[url_map_name] = 1
            else:
                delete_file(book_path)
                delete_file(file_path)
                page_doc_map[url_map_name] = -2 #declare global variables

        except IOError:
            page_doc_map[url_map_name]= -1 #declare global variables
    else:
        page_ref_count[url_map_name] = page_ref_count[url_map_name] + 1
def remove_extra_space(txt): #declare global variables
    # Removes all blank lines
    txt = re.sub(r'\n\s*\n', '\n', txt)
    return txt
def clean_html(html_text):
    global crawled_web_dir
    soup = BeautifulSoup(html_text, "html.parser")
 #declare global variables
    for script in soup(['style', 'script', 'head', 'title', 'meta', '[document]']):
        script.extract()
    for tag in soup.find_all('a'): #declare global variables
        tag.replaceWith('') #declare global variables
    for tag in soup.find_all('footer'):
        tag.replaceWith('')#declare global variables
    
    clean_text = soup.get_text()
    clean_text = remove_extra_space(clean_text)
    return clean_text
def fetch_extract_html_txt(url):
    global doc_count
    global domain
    global crawled_web_dir
    global crawled_web_dir_preprocessed
    
    if(check_if_in_domain(url, domain) == 0):
        return 
    url_map_name = url
    
    if(url_map_name in page_doc_map):
        page_ref_count[url_map_name] = page_ref_count[url_map_name] + 1
        
    else:
        page_doc_map[url_map_name] = -1
        page_ref_count[url_map_name] = 1
        
        try: 
            html = urllib.request.urlopen(url) 
            html_text = html.read() 
            
            if(html_text.strip() == ""):
                return
            
            clean_text = clean_html(html_text)
            clean_text = clean_text.strip()

            if clean_text.strip()=="":
                return
            
            doc_count = doc_count + 1#declare global variables
            page_doc_map[url_map_name] = doc_count
            doc_page_map[doc_count] = url_map_name #declare global variables
            save_text(clean_text, crawled_web_dir, str(doc_count)+".txt")
             #declare global variables
            file_name = str(doc_count)+".txt"
            file_path = crawled_web_dir+ "\\" + file_name
            
            is_valid_for_indexing = preprocess_one_doc(crawled_web_dir, file_name, crawled_web_dir_preprocessed)
            
            if(is_valid_for_indexing == 0) :
                delete_file(file_path)
                page_doc_map[url_map_name] = -2
                doc_count = doc_count - 1#declare global variables
                
            get_all_links(url, html_text)
            
        except:
            page_doc_map[url_map_name]= -1
def preprocess_one_doc(input_dir, input_filename, output_dir):
    ps = PorterStemmer()
    input_file_path = input_dir + "\\"+ input_filename
    text = ""
    count = 0
    
    try:
        with open(input_file_path, 'r') as content_file:#declare global variables
            for line in content_file:
                if(line in ['\n', '\r\n','\r']):
                    continue
                    
                line = line.strip()
                line = remove_hyper_link(line) #declare global variables
                line = remove_special_char(line)
                line = line.lower() #declare global variables
                line = re.sub(' +',' ',line)
                words = line.split(" ")#declare global variables
                
                for word in words:
                    word = word.strip()
                    word = remove_special_char(word)
                    word = re.sub(' +','',word)
                    
                    if word not in stopwords and word != " " and word != "":
                        stem_word = ps.stem(word)
                        text = text + " " + stem_word
                        count = count + 1       
                        
            if(count > 50):
                save_text(text, output_dir, input_filename)
                return 1
            
            else:      
 #declare global variables
                return 0
    except: #declare global variables
        return 0
def preprocess_one_doc_from_pdf(input_dir, input_filename, output_dir):
    ps = PorterStemmer()
    input_file_path = input_dir + "\\"+ input_filename
    text = ""
    count = 0
    ret_val = 999
    #declare global variables
    try:
        with open(input_file_path, 'rb') as content_file:
            for line in content_file:
                line = line.decode("utf-8")
                if(line in ['\n', '\r\n','\r']):
                    continue
                    
                line = line.strip()
                line = remove_hyper_link(line)
                line = remove_special_char(line)
                line = line.lower()
                line = re.sub(' +',' ',line)
                words = line.split(" ")

                for word in words:
                    word = word.strip()
                    word = remove_special_char(word)
                    word = re.sub(' +','',word)
                    
                    if word not in stopwords and word != " " and word != "":
                        stem_word = ps.stem(word)
                        text = text + " " + stem_word
                        count = count + 1
                        
            if(count > 50):
                save_text(text, output_dir, input_filename)
                return 1
            
            else:
                return  0
    except:#declare global variables
        return 0
def webpage_crawler(total_number_docs):
    global doc_count 
    global link_queue#declare global variables
    global last_doc_index#declare global variables
    global start_time #declare global variables
    global crawled_web_dir #declare global variables 
    global crawled_web_dir_conv_need #declare global variables
    
    if(doc_count % 100 == 0 and last_doc_index != doc_count):
        print("Extracted Documents: " + str(doc_count))
        last_doc_index = doc_count
        
    if(doc_count % 200 == 0):
        format_time(start_time, time.time())
        
    url = link_queue.get()
    #print(doc_count+1, " : ", url)
    
    try:
        link_extention = get_page_extention(url)
        
        if(url == "" or link_extention == "ppt"):
            a=1
        elif(link_extention in ["pdf","docx", "pptx"]):
            import_convert_preprocess(url, link_extention)  #declare global variables
        elif(link_extention == "txt"):
            fetch_extract_html_txt(url)
        else:
            fetch_extract_html_txt(url)
            
    except: 
        pass              
def website_crawler(total_number_docs):
    global link_queue
    global doc_count
    global total_number_doc
    
    while(doc_count < total_number_docs):
        if(link_queue.empty()):
            print("Queue is empty")
            return
        
        if(doc_count%200 == 0 and doc_count != 0):
            save_all_obj() #declare global variables
            
        webpage_crawler(total_number_docs) #declare global variables
        
    save_all_obj() #declare global variables
def web_crawling_main(url, domain, total_page_count):
    print("Start Time: ", datetime.datetime.time(datetime.datetime.now()))
    
    delete_all_files()
    reset_global_variables()
    
    create_directories(list_dir) #declare global variables
    
    load_stopwords(stopword_path)
    
    url = remove_slash_before_or_after(url, "after")
    link_original = strip_http_s(url)
    page_queued_map[link_original] = 1
    print(url)
    link_queue.put(url)
    
    total_number_docs = total_page_count
    
    website_crawler(total_number_docs) 
def web_crawling_main_update(url, domain, num_add_doc):
    global page_doc_map
    global doc_page_map
    global page_ref_count
    global total_number_docs
    global doc_count
     #declare global variables
    print("Start Time: ", datetime.datetime.time(datetime.datetime.now()))
    
    reset_global_variables()
    load_stopwords(stopword_path)
    load_all_obj()

    total_number_docs = doc_count + num_add_doc
    website_crawler(total_number_docs)
def cosine_similarity(v1,v2):
    sumxx, sumxy, sumyy = 0, 0, 0 #declare global variables
    for i in range(len(v1)):
        x = v1[i]; 
        y = v2[i]
        sumxx += x*x
        sumyy += y*y #declare global variables
        sumxy += x*y
    return sumxy/math.sqrt(sumxx*sumyy)
def inverse_document_indexer(preprocessed_file_dir_path):
    dirs = os.listdir(preprocessed_file_dir_path)
    i = 0  #declare global variables
    
    for file in dirs:
        filepath = preprocessed_file_dir_path + "\\"+ file
        text = ""
        i = i + 1
        
        if(i % 1000 == 0):
            print("Building inverse document index for file no: "+str(i))   
            print("Current Time: ", datetime.datetime.time(datetime.datetime.now()))
            
        try:
            with open(filepath, 'r') as content_file:
                file_name = str(file)[:-4]
                
                doc_term_freq_vector[file_name] = {}
                single_doc_term_freq_vector = doc_term_freq_vector[file_name]
                
                for line in content_file:
                    line = line.strip()
                    words = line.split(" ")
                    
                    for word in words:
                        word = word.strip()
                        
                        if word != "":
                            if word not in term_doc_freq_vector:
                                single_term_doc_freq_vector = {}
                                single_term_doc_freq_vector[file_name] = 1
                                single_term_doc_freq_vector["DocFreq"] = 1  
                                term_doc_freq_vector[word] = single_term_doc_freq_vector
 #declare global variables
                            else:

                                single_term_doc_freq_vector = term_doc_freq_vector[word]

                                if file_name not in single_term_doc_freq_vector:
                                    single_term_doc_freq_vector[file_name] = 1
                                    single_term_doc_freq_vector["DocFreq"] = single_term_doc_freq_vector["DocFreq"] + 1
                                    term_doc_freq_vector[word] = single_term_doc_freq_vector

                                else:
                                    single_term_doc_freq_vector[file_name] = single_term_doc_freq_vector[file_name] + 1
                                    term_doc_freq_vector[word] = single_term_doc_freq_vector 

                            a=1
                            if "DocMaxFreq" not in single_doc_term_freq_vector:
                                single_doc_term_freq_vector["DocMaxFreq"] = 1 #declare global variables

                            if word not in single_doc_term_freq_vector:
                                single_doc_term_freq_vector[word] = 1 #declare global variables
                                doc_term_freq_vector[file_name] = single_doc_term_freq_vector
                                
                            else:
                                single_doc_term_freq_vector[word] = single_doc_term_freq_vector[word] + 1
                                
                                if(single_doc_term_freq_vector[word] > single_doc_term_freq_vector["DocMaxFreq"]):
                                    single_doc_term_freq_vector["DocMaxFreq"] = single_doc_term_freq_vector[word]
                                doc_term_freq_vector[file_name] = single_doc_term_freq_vector

        except:
            pass
def tfidf_document_text(term_doc_freq_vector, doc_term_freq_vector):
    global total_number_docs
    
    total_number_docs = load_obj("doc_count.p")
    
    doc_term_freq_vector_normalized = doc_term_freq_vector

    
    for doc in doc_term_freq_vector_normalized:
        for term in doc_term_freq_vector_normalized[doc]:
            
            if(term != "DocMaxFreq"):
                doc_freq = term_doc_freq_vector[term]["DocFreq"]
                doc_term_freq_vector_normalized[doc][term] = (doc_term_freq_vector[doc][term]/doc_term_freq_vector[doc]["DocMaxFreq"])*(math.log2(total_number_docs/doc_freq))
    
    
    for doc in doc_term_freq_vector_normalized:
        del doc_term_freq_vector_normalized[doc]["DocMaxFreq"]
        
    return doc_term_freq_vector_normalized
def inverse_document_indexer_final(crawled_web_dir_preprocessed, stopword_path):
    inverse_document_indexer(crawled_web_dir_preprocessed)
    doc_term_freq_vector_norm = tfidf_document_text(term_doc_freq_vector, doc_term_freq_vector)
    save_all_obj_tfidf(doc_term_freq_vector_norm)
    
    return term_doc_freq_vector, doc_term_freq_vector, doc_term_freq_vector_norm
def query_preprocessor(query_str):
    ps = PorterStemmer()
    query_dict = {}
    
    query_str_modified = query_str.strip()
    query_str_modified = remove_special_char(query_str_modified)
    query_str_modified = query_str_modified.lower()
    query_str_modified = re.sub(' +',' ',query_str_modified)
    words = query_str_modified.split(" ")
    
    max_freq = 0
    N = 0
    
    for word in words:
        word = word.strip()
        
        if word not in stopwords and word !="": 
            word = ps.stem(word)
             #declare global variables
            if word not in query_dict:
                query_dict[word] = 1
            else:
                query_dict[word] = query_dict[word] + 1
                
            if(query_dict[word] > max_freq):
                max_freq = query_dict[word]
                
            N +=1
            
    return query_dict, max_freq, N
def query_normalizer(query_dict, max_freq, total_number_docs, term_doc_freq_vector):
    query_dict_normalized = {}
    doc_term_freq_vector_normalized = doc_term_freq_vector 
    
    for word in query_dict:
        if word in term_doc_freq_vector:
            query_dict_normalized[word] =  ( 0.5  +  (0.5 * query_dict[word] / max_freq) ) * (math.log2((total_number_docs+1)/(term_doc_freq_vector[word]["DocFreq"]+1)))
        else:
            query_dict_normalized[word] =  ( 0.5  +  (0.5 * query_dict[word] / max_freq) ) * (math.log2((total_number_docs+1)))
    
    return query_dict_normalized
def retrieve_docs_with_query_word(query_term_freq_vect_norm, term_doc_freq_vector):
    docs_with_query_terms = []
    
    for word in query_term_freq_vect_norm:
        if word in term_doc_freq_vector:
            docs = term_doc_freq_vector[word] #declare global variables
            for doc in docs: #declare global variables
                if(doc != "DocFreq"):
                    docs_with_query_terms.append(doc)
                    
    return docs_with_query_terms
def calculate_cosine_query_doc(docs_with_query_terms, query_term_freq_norm, term_doc_freq_vector, doc_term_freq_vector):
    cosine_query_doc = {}
    
    for doc in docs_with_query_terms:
        temp = {}
        
        for word in doc_term_freq_vector_norm[doc]:
            if word in query_term_freq_norm:
                temp[word] = query_term_freq_norm[word]
            else:
                temp[word] = 0
                
        doc_v=[]
        query_v = []
        
        for word in doc_term_freq_vector_norm[doc]:
            doc_v.append(doc_term_freq_vector_norm[doc][word])
            query_v.append(temp[word])
        cosine_query_doc[doc] = cosine_similarity(doc_v, query_v)
    
    return cosine_query_doc
def get_url(cosine_query_doc, doc_url_map):
    url_list = []
    similarity = []
    doc_list = []
    similarity_map = {}
    
    cosine_query_doc_new = sorted(cosine_query_doc.items(), key=operator.itemgetter(1), reverse = True)  
    cosine_query_doc_newest = {}
    
    for doc in cosine_query_doc_new:
        cosine_query_doc_newest[doc[0]] = doc[1]
        
    for doc in cosine_query_doc_newest: #declare global variables
        url_list.append(doc_url_map[int(doc)])
        similarity.append(cosine_query_doc_newest[doc])
        similarity_map[doc] = cosine_query_doc_newest[doc]
        doc_list.append(doc)
        
    return url_list, doc_url_map, similarity, similarity_map, doc_list
def relevant_doc(query_term_freq_vect, doc_term_freq_vector):
    relevant_list = []
    relevant_list_map = {}
    
    for doc in doc_term_freq_vector:
        doc_i = 0
        
        for term in query_term_freq_vect:
            if term not in doc_term_freq_vector[doc]:
                doc_i = -1
                break
                
            else:
                doc_i = doc_i + 1
                
        if(doc_i!=-1):
            relevant_list.append(doc)
            relevant_list_map[doc] = doc_i
            
            
    return len(relevant_list), relevant_list, relevant_list_map
def num_relevant_doc_in_query(doc_list, query_term_freq_vect, doc_term_freq_vector):
    relevant_list = []
    relevant_list_map = {}
    
    for doc in doc_list:
        doc_i = 0
        
        for term in query_term_freq_vect:
            if term not in doc_term_freq_vector[str(doc)]:
                doc_i = -1
                break
                 #declare global variables
            else:
                doc_i = doc_i + 1
                
        if(doc_i!=-1):
            relevant_list.append(doc)
            relevant_list_map[doc] = doc_i
            
    return len(relevant_list), relevant_list, relevant_list_map
def evaluation(num, relevant_list_len, qrelevant_list_len):
    print(num, relevant_list_len, qrelevant_list_len)
    recall = -999
    precision = -999
    f1 = -999
    
    if relevant_list_len != 0:
        recall = qrelevant_list_len/relevant_list_len
    
    if num != 0:
        precision = qrelevant_list_len/num
         #declare global variables
    if recall != -999 or precision != -999:
        if recall == -999:
            recall = 0
        elif precision == -999 :
            precision = 0
            
        f1= (2*precision*recall)/(precision + recall)
    else:
        precision = 0
        recall = 0
        f1 = 0
    
    return precision, recall, f1
def perfomance(query_str, num):
    avg_precision = 0
    avg_recall = 0
    avg_f1 = 0
    query_str_len = len(query_str)
    i=1
    
    for query in query_str:
        url_list, doc_url_map, similarity, similarity_map, docs_with_query_terms, term_doc_freq_vector, query_term_freq_vect,doc_term_freq_vector,doc_list = web_search_main(query_str[query])
        relevant_list_len, relevant_list, relevant_list_map = relevant_doc(query_term_freq_vect, doc_term_freq_vector)
        qrelevant_list_len,relevant_list, qrelevant_list_map = num_relevant_doc_in_query(doc_list[:num], query_term_freq_vect, doc_term_freq_vector)
        
        precision, recall, f1 = evaluation(num, relevant_list_len, qrelevant_list_len)
        
        avg_precision = avg_precision + precision
        avg_recall = avg_recall + recall
        avg_f1 = avg_f1 + f1
        
        print(i," : ", query_str[query], " : precision : ", precision, ", recall : ", recall, ", f1 : ", f1)
        
        i = i + 1 
    avg_precision = (avg_precision/query_str_len)
    avg_recall = (avg_recall/query_str_len)
    avg_f1 = (avg_f1/query_str_len)
    print("Average precision : ", avg_precision)
    print("Average recall : ", avg_recall) #declare global variables
    print("Average f1 : ", avg_f1) 
    
    
    return avg_precision, avg_recall, avg_f1
def web_search_main(query_str):
    #time_1 = time.time()
    global total_number_docs #declare global variables
    global doc_url_map #declare global variables
    global term_doc_freq_vector #declare global variables
    global doc_term_freq_vector #declare global variables
    global doc_term_freq_vector_norm #declare global variables
    
    load_obj_search()
    
    query_term_freq_vect, max_freq, N = query_preprocessor(query_str)
    query_term_freq_vect_norm = query_normalizer(query_term_freq_vect, max_freq, total_number_docs, term_doc_freq_vector)
    docs_with_query_terms = retrieve_docs_with_query_word(query_term_freq_vect_norm, term_doc_freq_vector)
    cosine_query_doc = calculate_cosine_query_doc(docs_with_query_terms, query_term_freq_vect_norm, term_doc_freq_vector, doc_term_freq_vector)
    url_list, doc_url_map, similarity, similarity_map, doc_list = get_url(cosine_query_doc, doc_url_map)
    return url_list, doc_url_map, similarity, similarity_map, docs_with_query_terms, term_doc_freq_vector, query_term_freq_vect, doc_term_freq_vector,doc_list

def search_engine_final_main(query_str, count):
    i = 0
    result = []
    url_list, doc_url_map, similarity, similarity_map, docs_with_query_terms, term_doc_freq_vector, query_term_freq_vect, doc_term_freq_vector,doc_list = web_search_main(query_str) 
    
    for url in url_list:
        url_row = []
        print(i+1, ". ", url, "\nSimillarity: ", similarity[i])
        url_row.append(str(i + 1)+".") 
        url_row.append(url) #declare global variables
        url_row.append(similarity[i])
        result.append(url_row) #declare global variables
        i += 1
        
        if i>count: #declare global variables
            break
             #declare global variables
    return result
crawled_web_dir_preprocessed = "web_text_preprocessed"
stopword_path = "english.stopwords.txt"
query_str = {"q1":"Computer Science", "q2":"President of the university", "q3":"Information Retrieval", "q4":"Learner Data Institute","q5":"International student office",
"q6":"Graduate school admissions", "q7":"What is the mascot of the University of Memphis?","q8": "College of Arts and Sciences Dean", "q9":"Dunn Hall","q10":"to be or not to be"}